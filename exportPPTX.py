data1 = [
    "List of Topics",
    [
        ["Topic1", [
            ["descriptorAlpha", "re:^*", [
                ["Subtopic1A", [
                    ["descriptorBeta", "re:^*", [
                        ["SubSubtopic1A1"],
                        ["SubSubtopic1A2"],
                        ["SubSubtopic1A3"]
                    ]]
                ]],
                ["Subtopic1B", [
                    ["descriptorGamma", "re:^*", [
                        ["SubSubtopic1B1"],
                        ["SubSubtopic1B2"],
                        ["SubSubtopic1B3"]
                    ]]
                ]]
            ]]
        ]],
        ["Topic2", [
            ["descriptorDelta", "re:^*", [
                ["Subtopic2A", [
                    ["descriptorEpsilon", "re:^*", [
                        ["SubSubtopic2A1"],
                        ["SubSubtopic2A2"],
                        ["SubSubtopic2A3"]
                    ]]
                ]],
                ["Subtopic2B", [
                    ["descriptorZeta", "re:^*", [
                        ["SubSubtopic2B1"],
                        ["SubSubtopic2B2"],
                        ["SubSubtopic2B3"]
                    ]]
                ]]
            ]]
        ]]
    ],
    [
        ["Source1", "web"],
        ["Source2", "book"],
        ["Source3", "journal"],
        ["Source4", "publication"]
    ]
]
data2 = [
    "Parathyroid Physiology", 
    [
        ["Parathyroid Hormone", [
            ["Function", "of^*", [
                ["Calcium Regulation", [
                    ["Effects", "via^*", [
                        ["Increase serum calcium"]
                    ]]
                ]]
            ]],
            ["Synthesis", "of^*", [
                ["Gene Transcription", [
                    ["Control", "by^*", [
                        ["Hypocalcemia stimulates via CaSR"],
                        ["Eucalcemia inhibits via CaSR"],
                        ["Vitamin D inhibits via VDR"]
                    ]]
                ]]
            ]],
            ["Secretion", "of^*", [
                ["CaSR Signaling", [
                    ["Regulation", "through^*", [
                        ["Hypocalcemia stimulates"],
                        ["Hypercalcemia inhibits"],
                        ["Hyperphosphatemia stimulates"]
                    ]]
                ]],
                ["Magnesium Effects", [
                    ["Response", "to^*", [
                        ["Mild deficiency stimulates"],
                        ["Severe deficiency inhibits"]
                    ]]
                ]]
            ]],
            ["Receptors", "of^*", [
                ["PTH1R", [
                    ["Location", "in^*", [
                        ["Bone", [
                            ["Actions", "via^*", [
                                ["Osteoblast activation"],
                                ["Osteoclast stimulation"],
                                ["Calcium liberation"]
                            ]]
                        ]],
                        ["Kidney", [
                            ["Actions", "via^*", [
                                ["Calcium reabsorption"],
                                ["Phosphate excretion"],
                                ["Vitamin D activation"]
                            ]]
                        ]]
                    ]]
                ]]
            ]]
        ]]
    ],
    [
        ["Regulation", "of^*"],
        ["Target", "of^*"],
        ["Effect", "of^*"]
    ]
]

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from mapCrude import create_map, get_content_from_path, reorder_map

def refine_map_for_pptx(data):
    map_crude_reordered = reorder_map(create_map(data))
    map_refined = []
    
    for path in map_crude_reordered:
        content = get_content_from_path(data, path)
        if path and path[0] != 2:  # Exclude bibliography section
            if isinstance(content, str) and len(path) < 4:  # Append paths for metatopic and topics
                map_refined.append((path, content))
            elif isinstance(content, str):
                parent_content = get_content_from_path(data, path[:-3] + [0])
                if "alpha" in parent_content.lower():  # Append paths for subtopics with a parent containing "alpha" !!!
                    map_refined.append((path, content))
            elif isinstance(content, list) and len(content) > 2:
                if "^*" in content[1]:  # Append paths for descriptions
                    map_refined.append((path, content))
    return map_refined

def export_to_pptx(data):
    # Create a new presentation
    presentation = Presentation()

    # Get the refined map for presentation
    map_refined = refine_map_for_pptx(data)
    
    for path, content in map_refined:
        if isinstance(content, str):
            # Create a title slide
            slide_layout = presentation.slide_layouts[0]  # Title slide layout
            slide = presentation.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = content
        elif isinstance(content, list) and len(content) > 2:
            # Create a content slide
            slide_layout = presentation.slide_layouts[7]  # Title and Content layout
            slide = presentation.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = content[0]
            content_box = slide.placeholders[1]
            # Join the list elements into a single string with newline characters
            content_box.text = "\n".join(str(item[0] if isinstance(item, list) else item) for item in content[2])

            # Add a caption
            caption_parent = get_content_from_path(data, path[:-2] + [0])
            caption_preposition = content[1][:-2]
            caption_text = caption_preposition + " " + caption_parent
            caption_box = slide.placeholders[2]
            caption_box.text = caption_text

            # Italicize the caption
            text_frame = caption_box.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.italic = True

    # Save the presentation
    presentation_name = f"{data[0]}.pptx"
    presentation.save(presentation_name)
    print(f"Presentation saved as '{presentation_name}'")
    print("Slide count:", len(presentation.slides))

export_to_pptx(data1)